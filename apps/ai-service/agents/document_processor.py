"""Document Processing Agent - extracts text and structure from documents."""

from __future__ import annotations

import logging
import os
from typing import Any

from core.agent_base import AgentBase
from models.schemas import AgentMessage

logger = logging.getLogger(__name__)


class DocumentProcessorAgent(AgentBase):
    """Extracts text, tables, images, and metadata from document formats."""

    def __init__(self) -> None:
        super().__init__(
            agent_name="document_processor",
            description="Processes PDF, DOCX, PPTX, plain text, markdown, and HTML documents.",
        )

    async def handle(self, message: AgentMessage) -> dict[str, Any]:
        """Process a document file and return extracted content."""
        payload = message.payload
        file_path: str = payload.get("file_path", "")
        extension: str = payload.get("extension", "")

        if not file_path or not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}

        ext = extension or os.path.splitext(file_path)[1].lower()

        try:
            if ext == ".pdf":
                return await self._extract_pdf(file_path)
            elif ext == ".docx":
                return await self._extract_docx(file_path)
            elif ext == ".pptx":
                return await self._extract_pptx(file_path)
            elif ext in (".txt", ".log"):
                return await self._extract_plain_text(file_path)
            elif ext == ".md":
                return await self._extract_markdown(file_path)
            elif ext in (".html", ".htm"):
                return await self._extract_html(file_path)
            else:
                return await self._extract_plain_text(file_path)
        except Exception as exc:
            self.logger.error("Failed to process document %s: %s", file_path, exc)
            return {"error": str(exc), "file_path": file_path}

    async def _extract_pdf(self, file_path: str) -> dict[str, Any]:
        """Extract text, tables, images, and metadata from PDF using PyMuPDF."""
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        metadata = doc.metadata or {}

        pages: list[dict[str, Any]] = []
        full_text_parts: list[str] = []
        total_images = 0
        tables_found: list[dict[str, Any]] = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text
            text = page.get_text("text")
            full_text_parts.append(text)

            # Extract images
            image_list = page.get_images(full=True)
            total_images += len(image_list)

            # Try table extraction
            try:
                page_tables = page.find_tables()
                for table in page_tables:
                    table_data = table.extract()
                    if table_data:
                        tables_found.append({
                            "page": page_num + 1,
                            "rows": len(table_data),
                            "cols": len(table_data[0]) if table_data else 0,
                            "data": table_data[:10],  # Preview first 10 rows
                        })
            except Exception:
                pass  # Table extraction not available in all PyMuPDF versions

            pages.append({
                "page_number": page_num + 1,
                "text_length": len(text),
                "image_count": len(image_list),
                "text_preview": text[:500] if text else "",
            })

        doc.close()

        full_text = "\n\n".join(full_text_parts)

        return {
            "format": "pdf",
            "page_count": len(pages),
            "metadata": {
                "title": metadata.get("title", ""),
                "author": metadata.get("author", ""),
                "subject": metadata.get("subject", ""),
                "creator": metadata.get("creator", ""),
                "creation_date": metadata.get("creationDate", ""),
            },
            "full_text": full_text,
            "text_length": len(full_text),
            "pages": pages,
            "total_images": total_images,
            "tables": tables_found,
            "file_path": file_path,
        }

    async def _extract_docx(self, file_path: str) -> dict[str, Any]:
        """Extract text and structure from DOCX using python-docx."""
        from docx import Document

        doc = Document(file_path)

        paragraphs: list[dict[str, Any]] = []
        full_text_parts: list[str] = []
        tables_found: list[dict[str, Any]] = []
        headings: list[dict[str, str]] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else "Normal"
            full_text_parts.append(text)

            para_info: dict[str, Any] = {
                "text": text,
                "style": style_name,
            }

            if style_name.startswith("Heading"):
                headings.append({"level": style_name, "text": text})

            paragraphs.append(para_info)

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            rows_data: list[list[str]] = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_data)

            tables_found.append({
                "table_index": table_idx,
                "rows": len(rows_data),
                "cols": len(rows_data[0]) if rows_data else 0,
                "data": rows_data[:10],
            })

        full_text = "\n".join(full_text_parts)

        # Extract core properties
        props = doc.core_properties
        metadata = {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
        }

        return {
            "format": "docx",
            "metadata": metadata,
            "full_text": full_text,
            "text_length": len(full_text),
            "paragraph_count": len(paragraphs),
            "headings": headings,
            "tables": tables_found,
            "file_path": file_path,
        }

    async def _extract_pptx(self, file_path: str) -> dict[str, Any]:
        """Extract slides, notes, and images from PPTX using python-pptx."""
        from pptx import Presentation

        prs = Presentation(file_path)

        slides: list[dict[str, Any]] = []
        full_text_parts: list[str] = []
        total_images = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts: list[str] = []
            image_count = 0

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)
                            full_text_parts.append(text)

                if shape.shape_type == 13:  # Picture
                    image_count += 1
                    total_images += 1

            # Extract notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    full_text_parts.append(f"[Notes] {notes_text}")

            slides.append({
                "slide_number": slide_idx + 1,
                "text": "\n".join(slide_text_parts),
                "notes": notes_text,
                "image_count": image_count,
            })

        full_text = "\n\n".join(full_text_parts)

        return {
            "format": "pptx",
            "slide_count": len(slides),
            "slides": slides,
            "full_text": full_text,
            "text_length": len(full_text),
            "total_images": total_images,
            "file_path": file_path,
        }

    async def _extract_plain_text(self, file_path: str) -> dict[str, Any]:
        """Extract content from plain text files."""
        encoding = "utf-8"
        text = ""

        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read()
                encoding = enc
                break
            except (UnicodeDecodeError, ValueError):
                continue

        lines = text.split("\n")

        return {
            "format": "text",
            "encoding": encoding,
            "full_text": text,
            "text_length": len(text),
            "line_count": len(lines),
            "file_path": file_path,
        }

    async def _extract_markdown(self, file_path: str) -> dict[str, Any]:
        """Extract content from markdown files with structural elements."""
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        lines = text.split("\n")
        headings: list[dict[str, str]] = []

        for line in lines:
            stripped = line.strip()
            if stripped.startswith("#"):
                level = len(stripped) - len(stripped.lstrip("#"))
                heading_text = stripped.lstrip("#").strip()
                if heading_text:
                    headings.append({"level": f"h{level}", "text": heading_text})

        return {
            "format": "markdown",
            "full_text": text,
            "text_length": len(text),
            "line_count": len(lines),
            "headings": headings,
            "file_path": file_path,
        }

    async def _extract_html(self, file_path: str) -> dict[str, Any]:
        """Extract text from HTML, stripping tags."""
        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()

        # Simple tag stripping (for full HTML parsing, use beautifulsoup)
        import re

        text = re.sub(r"<script[^>]*>.*?</script>", "", html_content, flags=re.DOTALL)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()

        return {
            "format": "html",
            "full_text": text,
            "text_length": len(text),
            "raw_html_length": len(html_content),
            "file_path": file_path,
        }
