"""Document Q&A Agent - RAG-based question answering over documents via cloud LLM."""

from __future__ import annotations

import asyncio
import csv
import json
import logging
import os
import sqlite3
import subprocess
from typing import Any

from core.agent_base import AgentBase
from models.schemas import AgentMessage
from services.ollama_client import chat_completion, generate_embedding

logger = logging.getLogger(__name__)

DATABASE_PATH = os.getenv("DATABASE_PATH", "./data/databases")
UPLOAD_PATH = os.getenv("UPLOAD_PATH", "./data/uploads")

QA_SYSTEM = """You are a helpful data assistant for the Data Ruler platform.
Answer questions based on the provided context. If the context doesn't contain
enough information, say so honestly. Use markdown formatting.

When discussing data:
- Reference specific column names and values
- Suggest SQL queries when appropriate
- Recommend visualizations for data insights
- Be precise about numbers and statistics"""


class DocumentQAAgent(AgentBase):
    """RAG-based Q&A over documents and data using cloud LLM + embeddings."""

    def __init__(self) -> None:
        super().__init__(
            agent_name="document_qa",
            description="Answers questions about documents and data using RAG with cloud LLM and embeddings.",
        )

    async def handle(self, message: AgentMessage) -> dict[str, Any]:
        payload = message.payload
        question = payload.get("message", payload.get("question", ""))
        schema_context = payload.get("schema_context", "")
        document_context = payload.get("document_context", "")
        conversation_history = payload.get("conversation_history", [])
        user_id = payload.get("user_id", "")
        file_id = payload.get("file_id", "")

        if not question:
            return {"error": "No question provided"}

        # Try to extract document text for non-tabular files (PDF, DOCX, etc.)
        if not document_context and user_id:
            document_context = await self._fetch_document_text(user_id, file_id, question)

        # Build context from available sources
        context_parts = []
        if schema_context:
            context_parts.append(f"Available data tables:\n{schema_context}")
        if document_context:
            context_parts.append(f"Document content:\n{document_context}")

        full_context = "\n\n".join(context_parts) if context_parts else "No specific data context available."

        # Build conversation — keep recent history short to avoid exceeding context limits
        messages = []
        total_chars = 0
        for h in reversed(conversation_history[-10:]):
            content = h.get("content", "") or ""
            # Truncate individual messages and cap total history size
            if len(content) > 1500:
                content = content[:1500] + "..."
            if total_chars + len(content) > 6000:
                break
            messages.insert(0, {
                "role": h.get("role", "user"),
                "content": content,
            })
            total_chars += len(content)
        messages.append({
            "role": "user",
            "content": f"Context:\n{full_context}\n\nQuestion: {question}",
        })

        locale = payload.get("locale", "en")
        system = QA_SYSTEM
        if locale == "he":
            system += "\n\nIMPORTANT: Always respond entirely in Hebrew (עברית). All text, headings, and explanations must be in Hebrew."

        try:
            answer = await chat_completion(
                messages=messages,
                system=system,
                temperature=0.5,
                max_tokens=1500,
                model_tier="chat",
            )
            return {
                "answer": answer,
                "question": question,
                "context_used": bool(context_parts),
                "status": "success",
            }
        except Exception as exc:
            return {
                "error": str(exc),
                "question": question,
                "status": "error",
            }

    async def _fetch_document_text(
        self, user_id: str, file_id: str | None, question: str,
    ) -> str:
        """Look up document files in the catalog and extract text content.

        For non-tabular files (PDF, DOCX, TXT, etc.) the schema_context is
        empty or has 0 columns.  This method finds the actual file on disk,
        extracts text, and returns it (truncated to fit context limits).
        """
        try:
            return await asyncio.to_thread(
                self._fetch_document_text_sync, user_id, file_id, question,
            )
        except Exception as exc:
            logger.warning("Document text fetch failed: %s", exc)
            return ""

    @staticmethod
    def _fetch_document_text_sync(
        user_id: str, file_id: str | None, question: str,
    ) -> str:
        catalog_db = os.path.join(DATABASE_PATH, "catalog.db")
        if not os.path.exists(catalog_db):
            return ""

        conn = sqlite3.connect(catalog_db)
        conn.row_factory = sqlite3.Row
        try:
            if file_id:
                rows = conn.execute(
                    """SELECT id, original_name, stored_path, file_type,
                              file_category, mime_type, size_bytes, row_count,
                              column_count, schema_snapshot, quality_score,
                              quality_profile, ai_summary, tags
                       FROM files WHERE user_id = ? AND id = ?""",
                    (user_id, file_id),
                ).fetchall()
            else:
                rows = conn.execute(
                    """SELECT id, original_name, stored_path, file_type,
                              file_category, mime_type, size_bytes, row_count,
                              column_count, schema_snapshot, quality_score,
                              quality_profile, ai_summary, tags
                       FROM files WHERE user_id = ?
                       ORDER BY created_at DESC LIMIT 20""",
                    (user_id,),
                ).fetchall()

            # Match files mentioned in the question
            target_rows = []
            q_lower = question.lower()
            for row in rows:
                name = (row["original_name"] or "").lower()
                if name and name in q_lower:
                    target_rows.append(row)
            if not target_rows:
                target_rows = rows[:5]

            text_parts = []
            for row in target_rows:
                path = row["stored_path"]
                name = row["original_name"] or "unknown"

                # Always include file metadata — works for ANY file type
                meta_lines = [f"File: {name}"]
                if row["file_type"]:
                    meta_lines.append(f"Type: {row['file_type']}")
                if row["file_category"]:
                    meta_lines.append(f"Category: {row['file_category']}")
                if row["mime_type"]:
                    meta_lines.append(f"MIME: {row['mime_type']}")
                if row["size_bytes"]:
                    size_mb = row["size_bytes"] / (1024 * 1024)
                    meta_lines.append(f"Size: {size_mb:.1f} MB")
                if row["row_count"]:
                    meta_lines.append(f"Rows: {row['row_count']}")
                if row["column_count"]:
                    meta_lines.append(f"Columns: {row['column_count']}")
                if row["quality_score"] is not None:
                    meta_lines.append(f"Quality score: {row['quality_score']}")
                if row["tags"]:
                    meta_lines.append(f"Tags: {row['tags']}")

                # Schema info
                if row["schema_snapshot"]:
                    try:
                        schema = json.loads(row["schema_snapshot"])
                        if isinstance(schema, list) and schema:
                            cols = ", ".join(
                                f"{c.get('name', '?')} ({c.get('inferred_type', 'text')})"
                                for c in schema[:20]
                            )
                            meta_lines.append(f"Schema: {cols}")
                    except (json.JSONDecodeError, TypeError):
                        pass

                # Quality profile
                if row["quality_profile"]:
                    try:
                        profile = json.loads(row["quality_profile"])
                        if isinstance(profile, dict) and profile.get("issues"):
                            meta_lines.append(f"Quality issues: {', '.join(str(i) for i in profile['issues'][:5])}")
                    except (json.JSONDecodeError, TypeError):
                        pass

                # AI summary if available
                if row["ai_summary"]:
                    meta_lines.append(f"AI Summary: {row['ai_summary']}")

                # Try extracting text content for supported formats.
                # AI best practice: present data as structured markdown
                # tables/sections so the LLM can reason about it.
                content = ""
                if path and os.path.exists(path):
                    ext = os.path.splitext(path)[1].lower()
                    try:
                        if ext == ".pdf":
                            import fitz
                            doc = fitz.open(path)
                            pages_text = [p.get_text() for p in doc]
                            doc.close()
                            content = "\n\n".join(pages_text)[:6000]

                        elif ext in (".xlsx", ".xls"):
                            content = DocumentQAAgent._extract_excel_for_ai(path)

                        elif ext == ".docx":
                            content = DocumentQAAgent._extract_docx_for_ai(path)

                        elif ext == ".doc":
                            content = DocumentQAAgent._extract_doc_for_ai(path)

                        elif ext in (".pptx", ".ppt"):
                            content = DocumentQAAgent._extract_pptx_for_ai(path)

                        elif ext in (".csv", ".tsv"):
                            content = DocumentQAAgent._extract_csv_for_ai(path, ext)

                        elif ext in (".db", ".sqlite", ".sqlite3"):
                            content = DocumentQAAgent._extract_sqlite_for_ai(path)

                        elif ext in (".mdb", ".accdb"):
                            content = DocumentQAAgent._extract_access_for_ai(path)

                        elif ext == ".sql":
                            with open(path, "r", encoding="utf-8", errors="replace") as f:
                                content = f.read(6000)

                        elif ext in (".txt", ".md", ".html", ".json", ".xml",
                                     ".yaml", ".yml", ".toml", ".ini", ".log"):
                            with open(path, "r", encoding="utf-8", errors="replace") as f:
                                content = f.read(6000)
                    except Exception as exc:
                        meta_lines.append(f"Text extraction error: {exc}")

                if content:
                    meta_lines.append(f"Content:\n{content}")

                text_parts.append("\n".join(meta_lines))

            return "\n\n---\n\n".join(text_parts)
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # AI-optimised content extractors
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_excel_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract ALL sheets from an Excel file as markdown tables.

        AI best practice: present tabular data as markdown so the LLM
        can see column headers alongside values, making it far easier
        to reason about the data than raw CSV or JSON.
        """
        import openpyxl

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        budget = max_chars

        for sheet_name in wb.sheetnames:
            if budget <= 0:
                parts.append(f"\n### Sheet: {sheet_name}\n(truncated — context limit)")
                continue

            ws = wb[sheet_name]
            rows_iter = ws.iter_rows(values_only=True)
            header_row = next(rows_iter, None)
            if not header_row:
                continue

            headers = [str(h) if h is not None else f"col_{i}"
                       for i, h in enumerate(header_row)]

            section = [f"\n### Sheet: {sheet_name} ({ws.max_row or '?'} rows, {len(headers)} cols)\n"]
            # Markdown table header
            section.append("| " + " | ".join(headers) + " |")
            section.append("| " + " | ".join("---" for _ in headers) + " |")

            row_count = 0
            for row_vals in rows_iter:
                cells = [str(v) if v is not None else "" for v in row_vals[:len(headers)]]
                line = "| " + " | ".join(c[:60] for c in cells) + " |"
                section.append(line)
                row_count += 1
                # Keep a reasonable number of sample rows per sheet
                if row_count >= 30:
                    section.append(f"| ... ({(ws.max_row or row_count) - row_count} more rows) |")
                    break

            sheet_text = "\n".join(section)
            if len(sheet_text) > budget:
                sheet_text = sheet_text[:budget] + "\n...(truncated)"
            parts.append(sheet_text)
            budget -= len(sheet_text)

        wb.close()
        return "\n".join(parts) if parts else ""

    @staticmethod
    def _extract_docx_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract text AND tables from a Word document.

        AI best practice: include structural elements (headings, tables)
        rather than flat paragraph text, so the LLM understands the
        document layout.
        """
        from docx import Document as DocxDocument

        doc = DocxDocument(path)
        parts: list[str] = []
        budget = max_chars

        # Extract paragraphs with style awareness
        for para in doc.paragraphs:
            if budget <= 0:
                break
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower()
            if "heading" in style:
                level = 1
                for ch in style:
                    if ch.isdigit():
                        level = int(ch)
                        break
                line = f"{'#' * level} {text}"
            else:
                line = text
            parts.append(line)
            budget -= len(line) + 1

        # Extract tables as markdown
        for idx, table in enumerate(doc.tables):
            if budget <= 100:
                break
            section = [f"\n**Table {idx + 1}:**\n"]
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip()[:40] for cell in row.cells]
                section.append("| " + " | ".join(cells) + " |")
                if r_idx == 0:
                    section.append("| " + " | ".join("---" for _ in cells) + " |")
                if r_idx >= 20:
                    section.append(f"| ... ({len(table.rows) - r_idx} more rows) |")
                    break
            table_text = "\n".join(section)
            if len(table_text) > budget:
                table_text = table_text[:budget] + "\n...(truncated)"
            parts.append(table_text)
            budget -= len(table_text)

        return "\n".join(parts) if parts else ""

    @staticmethod
    def _extract_doc_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract text from legacy .doc files using antiword.

        Antiword handles the old binary Word format reliably and
        outputs plain text that an LLM can reason about directly.
        """
        try:
            result = subprocess.run(
                ["antiword", path],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout:
                return result.stdout[:max_chars]
        except FileNotFoundError:
            logger.debug("antiword not installed — .doc extraction unavailable")
        except Exception as exc:
            logger.warning("antiword failed: %s", exc)
        # Fallback: try reading any embedded text
        try:
            with open(path, "rb") as f:
                raw = f.read(max_chars * 2)
            text = raw.decode("utf-8", errors="ignore")
            # Filter to printable text runs
            chunks = []
            current: list[str] = []
            for ch in text:
                if ch.isprintable() or ch in ("\n", "\t"):
                    current.append(ch)
                else:
                    if len(current) > 20:
                        chunks.append("".join(current))
                    current = []
            if len(current) > 20:
                chunks.append("".join(current))
            return "\n".join(chunks)[:max_chars] if chunks else ""
        except Exception:
            return ""

    @staticmethod
    def _extract_pptx_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract slides from PowerPoint as structured markdown.

        AI best practice: format each slide as a section with its
        number, title, body text, notes, and any tables.
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.debug("python-pptx not installed")
            return ""

        prs = Presentation(path)
        parts: list[str] = []
        budget = max_chars

        for slide_num, slide in enumerate(prs.slides, 1):
            if budget <= 0:
                parts.append(f"\n### Slide {slide_num}\n(truncated)")
                break

            section = [f"\n### Slide {slide_num}"]

            # Extract title
            if slide.shapes.title and slide.shapes.title.text:
                section.append(f"**{slide.shapes.title.text.strip()}**\n")

            # Extract all text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            section.append(text)

                # Extract tables
                if shape.has_table:
                    table = shape.table
                    section.append("")
                    for r_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip()[:40] for cell in row.cells]
                        section.append("| " + " | ".join(cells) + " |")
                        if r_idx == 0:
                            section.append("| " + " | ".join("---" for _ in cells) + " |")

            # Slide notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    section.append(f"\n_Notes: {notes[:200]}_")

            slide_text = "\n".join(section)
            if len(slide_text) > budget:
                slide_text = slide_text[:budget] + "\n...(truncated)"
            parts.append(slide_text)
            budget -= len(slide_text)

        total = len(prs.slides)
        header = f"**PowerPoint: {total} slides**\n"
        return header + "\n".join(parts)

    @staticmethod
    def _extract_csv_for_ai(
        path: str, ext: str = ".csv", max_chars: int = 6000,
    ) -> str:
        """Extract CSV/TSV as a markdown table with headers + sample rows.

        AI best practice: markdown tables let the LLM see each value
        aligned with its column header, dramatically improving data
        comprehension versus raw delimited text.
        """
        delimiter = "\t" if ext == ".tsv" else ","
        try:
            with open(path, "rb") as f:
                raw = f.read(8192)
            encoding = "utf-8"
            try:
                raw.decode("utf-8")
            except UnicodeDecodeError:
                encoding = "latin-1"

            with open(path, "r", encoding=encoding, errors="replace") as f:
                sample = f.read(8192)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=f"{delimiter};|\t")
                except csv.Error:
                    dialect = csv.excel
                    dialect.delimiter = delimiter

                reader = csv.reader(f, dialect)
                headers = next(reader, None)
                if not headers:
                    return ""

                md_parts = [
                    "| " + " | ".join(h[:30] for h in headers) + " |",
                    "| " + " | ".join("---" for _ in headers) + " |",
                ]
                budget = max_chars - len(md_parts[0]) - len(md_parts[1]) - 2
                row_count = 0
                total_rows = 0

                for row in reader:
                    total_rows += 1
                    if row_count < 50:
                        cells = [str(c)[:40] for c in row[:len(headers)]]
                        while len(cells) < len(headers):
                            cells.append("")
                        line = "| " + " | ".join(cells) + " |"
                        if budget - len(line) <= 0:
                            md_parts.append("| ... (more rows) |")
                            # Count remaining rows
                            for _ in reader:
                                total_rows += 1
                            break
                        md_parts.append(line)
                        budget -= len(line) + 1
                        row_count += 1
                    else:
                        pass  # just counting total_rows

                # Count remaining rows after early break
                for _ in reader:
                    total_rows += 1

                header = f"**CSV: {total_rows} rows, {len(headers)} columns**\n\n"
                return header + "\n".join(md_parts)
        except Exception as exc:
            logger.warning("CSV extraction failed: %s", exc)
            try:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read(max_chars)
            except Exception:
                return ""

    @staticmethod
    def _extract_sqlite_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract schema + sample data from SQLite/DB files.

        AI best practice: list every table with its schema and a few
        sample rows so the LLM can answer questions about the database
        structure and content.
        """
        try:
            conn = sqlite3.connect(path)
            conn.row_factory = sqlite3.Row
        except Exception as exc:
            logger.warning("SQLite open failed: %s", exc)
            return ""

        try:
            tables = conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' "
                "AND name NOT LIKE 'sqlite_%' ORDER BY name"
            ).fetchall()

            if not tables:
                return "Empty database (no tables)."

            parts = [f"**SQLite Database: {len(tables)} tables**\n"]
            budget = max_chars - len(parts[0])

            for tbl in tables:
                if budget <= 100:
                    parts.append(f"\n### Table: {tbl['name']}\n(truncated)")
                    break

                tbl_name = tbl["name"]
                cols_info = conn.execute(
                    f'PRAGMA table_info("{tbl_name}")'
                ).fetchall()
                col_names = [c["name"] for c in cols_info]
                col_types = [c["type"] or "TEXT" for c in cols_info]

                count = conn.execute(
                    f'SELECT COUNT(*) as c FROM "{tbl_name}"'
                ).fetchone()["c"]

                section = [
                    f"\n### Table: {tbl_name} ({count} rows, {len(col_names)} cols)",
                    "Schema: " + ", ".join(
                        f"{n} ({t})" for n, t in zip(col_names, col_types)
                    ),
                    "",
                    "| " + " | ".join(col_names[:15]) + " |",
                    "| " + " | ".join("---" for _ in col_names[:15]) + " |",
                ]

                sample_rows = conn.execute(
                    f'SELECT * FROM "{tbl_name}" LIMIT 10'
                ).fetchall()
                for row in sample_rows:
                    cells = [
                        str(row[c])[:40] if row[c] is not None else ""
                        for c in col_names[:15]
                    ]
                    section.append("| " + " | ".join(cells) + " |")
                if count > 10:
                    section.append(f"| ... ({count - 10} more rows) |")

                tbl_text = "\n".join(section)
                if len(tbl_text) > budget:
                    tbl_text = tbl_text[:budget] + "\n...(truncated)"
                parts.append(tbl_text)
                budget -= len(tbl_text)

            return "\n".join(parts)
        except Exception as exc:
            logger.warning("SQLite extraction failed: %s", exc)
            return ""
        finally:
            conn.close()

    @staticmethod
    def _extract_access_for_ai(path: str, max_chars: int = 6000) -> str:
        """Extract schema + sample data from MS Access (.mdb/.accdb) files.

        Uses mdbtools command-line utilities which handle both MDB and
        ACCDB formats.
        """
        try:
            result = subprocess.run(
                ["mdb-tables", "-1", path],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode != 0:
                return f"Access DB: could not list tables ({result.stderr.strip()})"

            tables = [t.strip() for t in result.stdout.strip().split("\n") if t.strip()]
            if not tables:
                return "Access DB: no tables found."

            parts = [f"**MS Access Database: {len(tables)} tables**\n"]
            budget = max_chars - len(parts[0])

            for tbl_name in tables:
                if budget <= 100:
                    parts.append(f"\n### Table: {tbl_name}\n(truncated)")
                    break

                section = [f"\n### Table: {tbl_name}"]

                # Schema via mdb-schema
                schema_result = subprocess.run(
                    ["mdb-schema", path, "--table", tbl_name],
                    capture_output=True, text=True, timeout=30,
                )
                if schema_result.returncode == 0 and schema_result.stdout:
                    for line in schema_result.stdout.split("\n"):
                        line = line.strip()
                        if line and not line.startswith("--"):
                            section.append(line)
                        if len("\n".join(section)) > 500:
                            break

                # Sample data via mdb-export (CSV output)
                export_result = subprocess.run(
                    ["mdb-export", path, tbl_name],
                    capture_output=True, text=True, timeout=30,
                )
                if export_result.returncode == 0 and export_result.stdout:
                    lines = export_result.stdout.strip().split("\n")
                    if lines:
                        headers = lines[0].split(",")
                        section.append("")
                        section.append(
                            "| " + " | ".join(h.strip('"')[:30] for h in headers) + " |"
                        )
                        section.append(
                            "| " + " | ".join("---" for _ in headers) + " |"
                        )
                        for row_line in lines[1:11]:
                            cells = row_line.split(",")
                            section.append(
                                "| " + " | ".join(
                                    c.strip('"')[:40] for c in cells[:len(headers)]
                                ) + " |"
                            )
                        if len(lines) > 11:
                            section.append(f"| ... ({len(lines) - 11} more rows) |")

                tbl_text = "\n".join(section)
                if len(tbl_text) > budget:
                    tbl_text = tbl_text[:budget] + "\n...(truncated)"
                parts.append(tbl_text)
                budget -= len(tbl_text)

            return "\n".join(parts)

        except FileNotFoundError:
            logger.debug("mdbtools not installed — Access DB extraction unavailable")
            return "Access database file detected. Install mdbtools for content extraction."
        except Exception as exc:
            logger.warning("Access extraction failed: %s", exc)
            return ""
