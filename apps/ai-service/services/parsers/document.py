"""Document parsers for PDF, DOCX, PPTX, and text files."""

import logging
import os
from typing import Any

logger = logging.getLogger(__name__)


def parse_pdf(file_path: str) -> dict[str, Any]:
    """Parse a PDF file using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text_parts = []
        metadata = {
            "page_count": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
        }

        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text}")

        doc.close()
        full_text = "\n\n".join(text_parts)
        return {
            "text": full_text,
            "metadata": metadata,
            "page_count": metadata["page_count"],
            "char_count": len(full_text),
        }
    except Exception as e:
        logger.error(f"PDF parse failed: {e}")
        return {"text": "", "error": str(e)}


def parse_docx(file_path: str) -> dict[str, Any]:
    """Parse a DOCX file."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)

        # Extract tables
        tables = []
        for table in doc.tables:
            headers = [cell.text for cell in table.rows[0].cells] if table.rows else []
            rows = []
            for row in table.rows[1:]:
                rows.append({headers[i]: cell.text for i, cell in enumerate(row.cells) if i < len(headers)})
            tables.append({"columns": headers, "rows": rows})

        return {
            "text": full_text,
            "tables": tables,
            "paragraph_count": len(paragraphs),
            "char_count": len(full_text),
        }
    except Exception as e:
        logger.error(f"DOCX parse failed: {e}")
        return {"text": "", "error": str(e)}


def parse_pptx(file_path: str) -> dict[str, Any]:
    """Parse a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        full_text_parts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            notes = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text

            text = "\n".join(slide_text)
            slides.append({
                "slide_number": slide_num + 1,
                "text": text,
                "notes": notes,
            })
            full_text_parts.append(f"--- Slide {slide_num + 1} ---\n{text}")

        full_text = "\n\n".join(full_text_parts)
        return {
            "text": full_text,
            "slides": slides,
            "slide_count": len(slides),
            "char_count": len(full_text),
        }
    except Exception as e:
        logger.error(f"PPTX parse failed: {e}")
        return {"text": "", "error": str(e)}


def parse_text(file_path: str) -> dict[str, Any]:
    """Parse a plain text file."""
    try:
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        text = raw.decode(encoding, errors="replace")
        return {
            "text": text,
            "encoding": encoding,
            "char_count": len(text),
            "line_count": text.count("\n") + 1,
        }
    except Exception as e:
        logger.error(f"Text parse failed: {e}")
        return {"text": "", "error": str(e)}


def parse_markdown(file_path: str) -> dict[str, Any]:
    """Parse a markdown file."""
    result = parse_text(file_path)
    result["format"] = "markdown"
    return result


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "doc": parse_docx,
    "pptx": parse_pptx,
    "ppt": parse_pptx,
    "txt": parse_text,
    "md": parse_markdown,
    "text": parse_text,
    "rtf": parse_text,
}


def parse_document(file_path: str, file_type: str) -> dict[str, Any]:
    """Parse any document file."""
    parser = PARSERS.get(file_type.lower(), parse_text)
    return parser(file_path)
